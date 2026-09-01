"""Spreadsheets, Word and PowerPoint, read as text.

Fetching HTML and text-layer PDFs and nothing else is a hole rather than a
limitation, because the citable numbers are so often the ones nobody published
as a web page. Statistical agencies, trade databases and procurement portals
put their actual figures in XLSX and their actual terms in DOCX, and a stack
that cannot read those cannot read its own primary sources.

Format is decided by SIGNATURE first, not by the Content-Type header or the URL
extension. Servers mislabel these constantly -- an HTML error page served as
application/vnd.ms-excel is common, and handing that to a parser produces
either a crash or, worse, nonsense that looks like data. OOXML files are zips
and legacy Office files are OLE2 compound documents, so the bytes settle it.

Legacy .doc and .ppt are detected and refused rather than guessed at: reading
them needs antiword/catdoc or LibreOffice, none of which is installed here.
Refusing by name is honest, and lets the caller fall through the rest of the
tier ladder. Note that .xls IS supported -- xlrd 2.x dropped xlsx and kept xls,
which is the opposite of what the version bump suggests.
"""
import csv
import io
import os

# OOXML (xlsx/docx/pptx) is a zip; legacy Office is an OLE2 compound file. Both
# are unambiguous, which is why they are checked before anything the server
# claims about the bytes.
ZIP_MAGIC = b"PK\x03\x04"
OLE2_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

# A national statistics export can run to six figures of rows. Nothing
# downstream reads that far and the extract cache would not hold it anyway.
MAX_ROWS = int(os.environ.get("DETHROTTLED_DOC_MAX_ROWS", "2000"))
MAX_CELL = int(os.environ.get("DETHROTTLED_DOC_MAX_CELL", "200"))

LEGACY = {"doc", "ppt"}

# The member whose presence identifies each OOXML format.
OOXML_MEMBERS = {
    "xl/workbook.xml": "xlsx",
    "word/document.xml": "docx",
    "ppt/presentation.xml": "pptx",
}

# What a zip's `mimetype` member declares. OpenDocument, EPUB and OOXML are all
# zips sharing the same four magic bytes, so the signature says "a zip" and
# nothing more -- but ODF and EPUB both put a `mimetype` member first and
# uncompressed, precisely so it can be read like this.
#
# ODF matters because it is what most European public bodies publish in, and
# an EPUB is a zip of XHTML. Both are ordinary things to be handed by a link.
MIMETYPE_MEMBERS = {
    "application/vnd.oasis.opendocument.text": "odt",
    "application/vnd.oasis.opendocument.spreadsheet": "ods",
    "application/vnd.oasis.opendocument.presentation": "odp",
    "application/epub+zip": "epub",
}

EXTENSIONS = {
    ".xlsx": "xlsx", ".xlsm": "xlsx", ".xls": "xls", ".csv": "csv",
    ".tsv": "csv", ".docx": "docx", ".pptx": "pptx", ".doc": "doc",
    ".ppt": "ppt", ".odt": "odt", ".ods": "ods", ".odp": "odp",
    ".epub": "epub", ".rtf": "rtf",
}

CONTENT_TYPES = {
    "spreadsheetml.sheet": "xlsx",
    "wordprocessingml.document": "docx",
    "presentationml.presentation": "pptx",
    "vnd.ms-excel": "xls",
    "msword": "doc",
    "vnd.ms-powerpoint": "ppt",
    "text/csv": "csv",
    "opendocument.text": "odt",
    "opendocument.spreadsheet": "ods",
    "opendocument.presentation": "odp",
    "epub+zip": "epub",
    "application/rtf": "rtf",
    "text/rtf": "rtf",
}


def _ooxml_kind(data: bytes) -> str:
    """Which zip-based document this is, from the members it contains.

    Reading the directory is cheap: it lives at the end of the file and needs
    no decompression.
    """
    import zipfile
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            names = set(archive.namelist())
            for member, kind in OOXML_MEMBERS.items():
                if member in names:
                    return kind
            if "mimetype" in names:
                declared = archive.read("mimetype").decode(
                    "ascii", "replace").strip()
                return MIMETYPE_MEMBERS.get(declared, "")
    except Exception:
        # A truncated or encrypted zip is not a document we can read, and
        # saying so beats raising out of a type probe.
        return ""
    return ""


def kind_of(data: bytes, url: str = "", content_type: str = "") -> str:
    """The format of these bytes, or "" if it is not a document we read.

    Signature first, deliberately. The header and the extension are consulted
    only for formats with no distinguishing magic, which is CSV and nothing
    else.
    """
    if not data:
        return ""

    claimed = (content_type or "").lower()
    suffix = os.path.splitext((url or "").split("?")[0].lower())[1]

    if data[:4] == ZIP_MAGIC:
        # A zip that is not OOXML is some other archive, which we do not read.
        return _ooxml_kind(data)

    if data[:8] == OLE2_MAGIC:
        # Legacy Office. Telling xls from doc from ppt needs the OLE2 directory,
        # so the claimed type decides. Being wrong is cheap here: a wrong guess
        # fails to parse rather than inventing data.
        for marker, kind in CONTENT_TYPES.items():
            if marker in claimed and kind in ("xls", "doc", "ppt"):
                return kind
        return EXTENSIONS.get(suffix, "xls")

    # RTF is plain text with a fixed opening token.
    if data[:5] == b"{\\rtf":
        return "rtf"

    # No signature, so only CSV qualifies -- and only when something says so,
    # because otherwise every HTML page looks like a one-column CSV.
    if "text/csv" in claimed or suffix in (".csv", ".tsv"):
        head = data[:2048].lstrip().lower()
        if head.startswith((b"<!doctype", b"<html")):
            return ""                                 # mislabelled HTML
        return "csv"
    return ""


def tidy_structured(value, limit: int = 0) -> str:
    """Normalise whitespace WITHOUT losing line breaks.

    fetch._tidy_text collapses every run of whitespace to one space, which is
    right for prose and destructive here. In a table the row boundary carries
    meaning: run two rows together and "Denmark | 2024 | 5120 Denmark | 2023 |
    4560" no longer says which number belongs to which year. Spreadsheets, OCR
    pages and slide decks all come through this instead.
    """
    lines = []
    for line in str(value or "").splitlines():
        line = " ".join(line.split())
        if line or (lines and lines[-1]):
            lines.append(line)
    out = "\n".join(lines).strip()
    if limit and len(out) > limit:
        out = out[:limit].rsplit(" ", 1)[0].strip()
    return out


def _cell(value) -> str:
    if value is None:
        return ""
    return str(value).replace("\n", " ").replace("|", "/").strip()[:MAX_CELL]


def _rows_to_text(title: str, rows, limit: int) -> list:
    """Rows as pipe-separated lines under a sheet heading.

    The header row is kept in place because a number is only citable if you can
    still see which column and which sheet it came from.
    """
    out, size = [], 0
    if title:
        out.append("## %s" % title)
        size += len(out[-1])
    for count, row in enumerate(rows):
        if count >= MAX_ROWS or size >= limit:
            break
        line = " | ".join(_cell(v) for v in row).strip(" |")
        if not line:
            continue
        out.append(line)
        size += len(line) + 1
    return out


def _from_xlsx(data: bytes, limit: int) -> str:
    import openpyxl
    book = openpyxl.load_workbook(io.BytesIO(data), read_only=True,
                                  data_only=True)
    try:
        parts, size = [], 0
        for sheet in book.worksheets:
            if size >= limit:
                break
            lines = _rows_to_text(sheet.title,
                                  sheet.iter_rows(values_only=True),
                                  limit - size)
            if len(lines) > 1:
                parts.extend(lines)
                size += sum(len(x) + 1 for x in lines)
        return "\n".join(parts)
    finally:
        book.close()


def _from_xls(data: bytes, limit: int) -> str:
    import xlrd
    book = xlrd.open_workbook(file_contents=data)
    parts, size = [], 0
    for sheet in book.sheets():
        if size >= limit:
            break
        rows = (sheet.row_values(i) for i in range(sheet.nrows))
        lines = _rows_to_text(sheet.name, rows, limit - size)
        if len(lines) > 1:
            parts.extend(lines)
            size += sum(len(x) + 1 for x in lines)
    return "\n".join(parts)


def _from_csv(data: bytes, limit: int) -> str:
    text = data.decode("utf-8", errors="replace")
    try:
        dialect = csv.Sniffer().sniff(text[:4096])
    except csv.Error:
        dialect = csv.excel
    rows = csv.reader(io.StringIO(text), dialect)
    return "\n".join(_rows_to_text("", rows, limit))


def _from_docx(data: bytes, limit: int) -> str:
    import docx
    document = docx.Document(io.BytesIO(data))
    parts, size = [], 0
    for para in document.paragraphs:
        if size >= limit:
            break
        text = para.text.strip()
        if text:
            parts.append(text)
            size += len(text) + 1
    # Tables carry the numbers in most reports, so they are not optional.
    for table in document.tables:
        if size >= limit:
            break
        rows = ([c.text for c in r.cells] for r in table.rows)
        lines = _rows_to_text("", rows, limit - size)
        parts.extend(lines)
        size += sum(len(x) + 1 for x in lines)
    return "\n".join(parts)


def _from_pptx(data: bytes, limit: int) -> str:
    import pptx
    deck = pptx.Presentation(io.BytesIO(data))
    parts, size = [], 0
    for number, slide in enumerate(deck.slides, 1):
        if size >= limit:
            break
        said = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                said.append(shape.text_frame.text.strip())
        if said:
            block = "## Slide %d\n%s" % (number, "\n".join(said))
            parts.append(block)
            size += len(block) + 1
    return "\n".join(parts)


def _from_odf(data: bytes, limit: int) -> str:
    """OpenDocument text, spreadsheet or presentation.

    One reader for all three, because odfpy exposes them the same way: the
    document is XML and its text lives in <text:p> and <text:h> elements
    wherever they appear, table cells included.

    Tables are read first so that a spreadsheet arrives as rows rather than as
    a run of loose values -- the same reason XLSX rows keep their line breaks.
    """
    from odf import teletype
    from odf import text as odftext
    from odf.opendocument import load
    from odf.table import TableRow

    document = load(io.BytesIO(data))
    lines, rows = [], 0

    for row in document.getElementsByType(TableRow):
        cells = [_cell(teletype.extractText(cell)) for cell in row.childNodes]
        cells = [c for c in cells if c]
        if cells:
            lines.append(" | ".join(cells))
            rows += 1
            if rows >= MAX_ROWS or sum(len(x) for x in lines) > limit:
                break

    if not lines:
        for element in (list(document.getElementsByType(odftext.H))
                        + list(document.getElementsByType(odftext.P))):
            got = teletype.extractText(element).strip()
            if got:
                lines.append(got)
            if sum(len(x) for x in lines) > limit:
                break

    return "\n".join(lines)


def _from_epub(data: bytes, limit: int) -> str:
    """An EPUB is a zip of XHTML, so it goes through the same extractor as a
    web page -- which is what those documents are."""
    import ebooklib
    from ebooklib import epub

    from . import extract as _fx

    book = epub.read_epub(io.BytesIO(data), options={"ignore_ncx": True})
    parts = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        got = _fx.extract(item.get_content().decode("utf-8", "replace"),
                          max_chars=limit)
        if got["ok"]:
            parts.append(got["text"])
        if sum(len(x) for x in parts) > limit:
            break
    return "\n\n".join(parts)


def _from_rtf(data: bytes, limit: int) -> str:
    """RTF via a pure-Python stripper. Not a rich conversion and does not need
    to be: the job is the words, not the formatting."""
    # Checked here as well as in kind_of(), because to_text() is public and a
    # caller naming the wrong kind would otherwise get its own bytes back as
    # "text" -- the stripper passes non-RTF through untouched.
    if not data[:5].startswith(b"{"):
        raise ValueError("not rtf")
    from striprtf.striprtf import rtf_to_text
    return rtf_to_text(data.decode("utf-8", "replace"), errors="ignore")


READERS = {"xlsx": _from_xlsx, "xls": _from_xls, "csv": _from_csv,
           "docx": _from_docx, "pptx": _from_pptx,
           "odt": _from_odf, "ods": _from_odf, "odp": _from_odf,
           "epub": _from_epub, "rtf": _from_rtf}


def to_text(data: bytes, kind: str, limit: int) -> tuple:
    """(text, reason). Empty text always arrives with a reason why."""
    if kind in LEGACY:
        # No local converter for OLE2 Word/PowerPoint. Named rather than
        # silently empty, so the caller can fall through the tier ladder and so
        # the gap stays visible if it ever starts mattering.
        return "", "legacy_binary_format:%s" % kind
    reader = READERS.get(kind)
    if reader is None:
        return "", "unsupported_document:%s" % (kind or "unknown")
    try:
        text = reader(data, limit)
    except Exception as exc:
        return "", "document_parse_failed:%s" % type(exc).__name__
    if not text.strip():
        return "", "document_empty:%s" % kind
    return text[:limit], ""
