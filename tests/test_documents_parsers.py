"""The document parsers, against files this test actually builds.

Format detection is covered elsewhere; this is the other half -- whether the
bytes of a real spreadsheet, document or deck come back as usable text. Every
file here is generated rather than checked in, so the tests exercise the same
libraries a user's install would and cannot drift from a stale fixture.

Each parser is skipped if its optional library is missing, because these live
behind the `documents` extra and an absent extra is not a failure.
"""
import io

import pytest

from dethrottled import documents as docs


def _missing(module):
    import importlib.util
    return importlib.util.find_spec(module) is None


# ── spreadsheets ─────────────────────────────────────────────────────────────

@pytest.mark.skipif(_missing("openpyxl"), reason="openpyxl not installed")
def test_xlsx_becomes_text():
    import openpyxl
    book = openpyxl.Workbook()
    sheet = book.active
    sheet.title = "Capacity"
    sheet.append(["country", "year", "megawatts"])
    sheet.append(["Denmark", 2024, 5120])
    sheet.append(["Denmark", 2023, 4560])
    buf = io.BytesIO()
    book.save(buf)

    text, reason = docs.to_text(buf.getvalue(), "xlsx", 8000)
    assert "Denmark" in text
    assert "5120" in text
    assert "Capacity" in text, "the sheet name is context worth keeping"


@pytest.mark.skipif(_missing("openpyxl"), reason="openpyxl not installed")
def test_xlsx_rows_stay_on_separate_lines():
    """The reason structured text does not go through the prose tidier: run two
    rows together and the table stops saying which number belongs to which
    year."""
    import openpyxl
    book = openpyxl.Workbook()
    book.active.append(["Denmark", 2024, 5120])
    book.active.append(["Denmark", 2023, 4560])
    buf = io.BytesIO()
    book.save(buf)
    text, _ = docs.to_text(buf.getvalue(), "xlsx", 8000)
    assert text.count("\n") >= 1


@pytest.mark.skipif(_missing("openpyxl"), reason="openpyxl not installed")
def test_row_cap_bounds_a_huge_sheet(monkeypatch):
    import openpyxl
    monkeypatch.setattr(docs, "MAX_ROWS", 5)
    book = openpyxl.Workbook()
    for i in range(200):
        book.active.append(["row", i])
    buf = io.BytesIO()
    book.save(buf)
    text, _ = docs.to_text(buf.getvalue(), "xlsx", 100000)
    assert text.count("\n") < 50, "a 200-row sheet must not arrive whole"


def test_csv_becomes_text():
    data = b"country,year,megawatts\nDenmark,2024,5120\nDenmark,2023,4560\n"
    text, _ = docs.to_text(data, "csv", 4000)
    assert "Denmark" in text and "5120" in text


def test_csv_with_odd_encoding_does_not_raise():
    """Servers serve latin-1 while claiming utf-8 all the time."""
    text, _ = docs.to_text("nom,valeur\ncafé,3\n".encode("latin-1"), "csv", 1000)
    assert "valeur" in text


# ── word and powerpoint ──────────────────────────────────────────────────────

@pytest.mark.skipif(_missing("docx"), reason="python-docx not installed")
def test_docx_becomes_text():
    import docx as python_docx
    document = python_docx.Document()
    document.add_heading("Tender notice", level=1)
    document.add_paragraph("Sealed bids are invited for the supply of "
                           "photovoltaic modules.")
    buf = io.BytesIO()
    document.save(buf)

    text, _ = docs.to_text(buf.getvalue(), "docx", 8000)
    assert "Tender notice" in text
    assert "photovoltaic" in text


@pytest.mark.skipif(_missing("pptx"), reason="python-pptx not installed")
def test_pptx_becomes_text():
    from pptx import Presentation
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "Quarterly capacity"
    slide.placeholders[1].text = "Installed base grew to 5120 megawatts."
    buf = io.BytesIO()
    deck.save(buf)

    text, _ = docs.to_text(buf.getvalue(), "pptx", 8000)
    assert "Quarterly capacity" in text
    assert "5120" in text


# ── failure behaviour ────────────────────────────────────────────────────────

def test_a_corrupt_file_gives_a_reason_not_an_exception():
    """Half a download is a normal thing to be handed."""
    text, reason = docs.to_text(b"PK\x03\x04 truncated rubbish", "xlsx", 1000)
    assert text == ""
    assert reason


def test_an_unknown_kind_is_refused_cleanly():
    text, reason = docs.to_text(b"anything", "wordperfect", 1000)
    assert text == ""
    assert reason


@pytest.mark.skipif(_missing("openpyxl"), reason="openpyxl not installed")
def test_the_character_limit_is_respected():
    import openpyxl
    book = openpyxl.Workbook()
    for i in range(500):
        book.active.append(["a much longer cell value here", i])
    buf = io.BytesIO()
    book.save(buf)
    text, _ = docs.to_text(buf.getvalue(), "xlsx", 300)
    assert len(text) <= 300


@pytest.mark.skipif(_missing("openpyxl"), reason="openpyxl not installed")
def test_cell_contents_are_bounded(monkeypatch):
    """One pathological cell must not become the whole extraction."""
    import openpyxl
    monkeypatch.setattr(docs, "MAX_CELL", 20)
    book = openpyxl.Workbook()
    book.active.append(["x" * 5000, "ok"])
    buf = io.BytesIO()
    book.save(buf)
    text, _ = docs.to_text(buf.getvalue(), "xlsx", 100000)
    assert "x" * 100 not in text


def test_cell_separator_cannot_be_forged():
    """A cell containing the column separator would otherwise invent columns
    that are not in the sheet."""
    assert "|" not in docs._cell("a|b")


def test_newlines_inside_a_cell_do_not_invent_rows():
    assert "\n" not in docs._cell("line one\nline two")
